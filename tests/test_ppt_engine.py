from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from drawing_yh.report.ppt_engine import (
    PptBuildConfig,
    _add_header,
    choose_layout,
    draft_bullets,
    draft_notes,
    load_draft_specs,
)


class FakeReportConfig:
    SLIDES = [("reports_draft_example", "Example", "Example", "draft")]

    @staticmethod
    def load_chunks(*, fix_img_for_pages=True):
        return {
            "reports_draft_example": """## Example

**PPT 页标题**:Example title
**PPT 副标题**:示例

**主图**
<img src="report_figs/example.png">

**右侧 3 bullets**
- 第一条
- 第二条

**讲者备注**
- 备注一
"""
        }


class PptEngineTests(unittest.TestCase):
    def test_shared_parser_reads_project_report_config(self):
        specs = load_draft_specs(FakeReportConfig)
        self.assertEqual(len(specs), 1)
        self.assertEqual(specs[0].title, "Example title")
        self.assertEqual(specs[0].subtitle, "示例")
        self.assertEqual(specs[0].images, ("report_figs/example.png",))
        self.assertEqual(draft_bullets(specs[0].body_md), ["第一条", "第二条"])
        self.assertIn("备注一", draft_notes(specs[0].body_md))

    def test_long_title_keeps_fixed_font_and_moves_body_down(self):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        config = PptBuildConfig(
            output=Path("out.pptx"),
            figure_dir=Path("figures"),
            template_candidates=(),
            title="Test",
        )
        title = (
            "Bulk-supported candidate B-cell communication axes: "
            "phagocytosis, antigen presentation, ECM, metabolism"
        )
        body_top = _add_header(slide, title, "B 参与 bulk-reversed 边", config)
        text_shapes = sorted(
            [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text],
            key=lambda shape: shape.top,
        )
        title_shape, subtitle_shape = text_shapes[:2]
        self.assertTrue(title_shape.text_frame.word_wrap)
        self.assertEqual(title_shape.text, title)
        self.assertEqual(title_shape.text_frame.paragraphs[0].runs[0].font.size.pt, 20)
        self.assertEqual(subtitle_shape.text_frame.paragraphs[0].runs[0].font.size.pt, 16)
        self.assertGreater(body_top, Inches(1.0))

    def test_candidate_layouts_cover_wide_and_semantic_double_images(self):
        with tempfile.TemporaryDirectory(prefix="ppt-engine-test-") as temp:
            root = Path(temp)
            wide = root / "wide.png"
            Image.new("RGB", (1600, 500), "white").save(wide)
            wide_plan = choose_layout([wide], ["一条较长的说明文字" * 4])
            self.assertEqual(wide_plan.name, "wide_bottom")

            first = root / "peptide_domain_VSIG4.png"
            second = root / "peptide_domain_LILRB5.png"
            Image.new("RGB", (1770, 922), "white").save(first)
            Image.new("RGB", (1967, 922), "white").save(second)
            paired = choose_layout(
                [first, second],
                [
                    "VSIG4: 12 条肽段均位于胞外域。",
                    "LILRB5: 14 条肽段均位于胞外域。",
                    "两个候选可由肽段拓扑区分胞内外。",
                ],
            )
            self.assertEqual(paired.name, "paired_summary", paired.scores)


if __name__ == "__main__":
    unittest.main()
