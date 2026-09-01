from __future__ import annotations

import importlib.util
import os
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Dict, Iterable, List, Mapping, Optional, Sequence, Set, Tuple

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

EMU_PER_IN = 914400
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BODY_TOP = Inches(1.0)
BODY_BOTTOM = SLIDE_H - Inches(0.25)
DIVIDER_COLOR = RGBColor(0xBE, 0x76, 0x00)


@dataclass(frozen=True)
class DraftSpec:
    slug: str
    heading: str
    title: str
    subtitle: str
    body_md: str
    images: Tuple[str, ...]


@dataclass(frozen=True)
class SectionSpec:
    name: str
    first_slug: str


@dataclass
class PptBuildConfig:
    output: Path
    figure_dir: Path
    template_candidates: Sequence[Path]
    title: str
    subtitle: str = ""
    figure_dirs: Sequence[Path] = field(default_factory=tuple)
    include_cover: bool = True
    auto_open: bool = True
    title_font_pt: int = 20
    subtitle_font_pt: int = 16
    compact_slugs: Set[str] = field(default_factory=set)
    sections: Sequence[SectionSpec] = field(default_factory=tuple)
    width_weights: Mapping[str, Sequence[float]] = field(default_factory=dict)


@dataclass(frozen=True)
class DraftSpec:
    slug: str
    heading: str
    title: str
    subtitle: str
    body_md: str
    images: Tuple[str, ...]


@dataclass(frozen=True)
class LayoutPlan:
    name: str
    score: float
    scores: Mapping[str, float]


def default_template_candidates() -> Tuple[Path, ...]:
    return (
        Path("D:/Projects/Retina_Aging/retina-bulk.pptx"),
        Path("D:/Projects/Retina_Aging/report/_archive_pptx/retina-bulk.pptx"),
        Path("D:/Projects/Retina_Aging/讨论/retina-bulk.pptx"),
    )


_FIELD_PATTERNS = {
    "title": re.compile(r"^\*\*PPT 页标题\*\*:(.+?)\s*$", re.MULTILINE),
    "subtitle": re.compile(r"^\*\*PPT 副标题\*\*:(.+?)\s*$", re.MULTILINE),
}
_IMAGE_RE = re.compile(r'<img\b[^>]*\bsrc="([^"]+)"', re.IGNORECASE)


def load_draft_specs(report_config: Any) -> List[DraftSpec]:
    """从项目 `_report_config.py` 的 SLIDES/load_chunks 读取统一汇报草稿。"""
    chunks = report_config.load_chunks(fix_img_for_pages=False)
    specs: List[DraftSpec] = []
    for slug, _, _, category in report_config.SLIDES:
        if category != "draft":
            continue
        chunk = chunks[slug]
        heading_match = re.search(r"^(#{2,3})\s+(.+?)\s*$", chunk, re.MULTILINE)
        fields: Dict[str, str] = {}
        for field_name, pattern in _FIELD_PATTERNS.items():
            matches = pattern.findall(chunk)
            if len(matches) != 1:
                raise ValueError(
                    f"{slug}: expected exactly one PPT {field_name}, found {len(matches)}"
                )
            fields[field_name] = matches[0].strip()
        body = chunk
        for pattern in _FIELD_PATTERNS.values():
            body = pattern.sub("", body)
        body = re.sub(r"\n{3,}", "\n\n", body).strip() + "\n"
        specs.append(
            DraftSpec(
                slug=slug,
                heading=heading_match.group(2).strip() if heading_match else "",
                title=fields["title"],
                subtitle=fields["subtitle"],
                body_md=body,
                images=tuple(_IMAGE_RE.findall(chunk)),
            )
        )
    return specs


def validate_draft_specs(
    report_config: Any,
    *,
    max_images: int = 2,
    max_subtitle_chars: int = 32,
) -> List[str]:
    errors: List[str] = []
    try:
        specs = load_draft_specs(report_config)
    except Exception as exc:
        return [str(exc)]
    for spec in specs:
        if not spec.title:
            errors.append(f"{spec.slug}: empty PPT title")
        if not spec.subtitle:
            errors.append(f"{spec.slug}: empty PPT subtitle")
        if len(spec.subtitle) > max_subtitle_chars:
            errors.append(
                f"{spec.slug}: PPT subtitle is too long ({len(spec.subtitle)} chars)"
            )
        if len(spec.images) > max_images:
            errors.append(
                f"{spec.slug}: too many images for one PPT slide ({len(spec.images)})"
            )
    return errors


def _resolve_template(candidates: Sequence[Path]) -> Path:
    for candidate in candidates:
        if Path(candidate).is_file():
            return Path(candidate)
    raise FileNotFoundError(f"PPT template not found: {[str(p) for p in candidates]}")


def _load_template(path: Path) -> Presentation:
    prs = Presentation(str(path))
    slide_ids = prs.slides._sldIdLst
    for slide_id in list(slide_ids):
        prs.part.drop_rel(slide_id.rId)
        slide_ids.remove(slide_id)
    return prs


def _plain_text(text: str) -> str:
    # Preserve escaped literal asterisks while stripping Markdown bold markers.
    # This keeps significance legends such as ``\*\* p<0.01`` intact in PPT.
    star_token = "\uf000"
    protected = text.replace(r"\*", star_token)
    return (
        re.sub(r"<[^>]+>", "", protected)
        .replace("**", "")
        .replace("`", "")
        .replace("&gt;", ">")
        .replace("&lt;", "<")
        .replace(star_token, "*")
        .strip()
    )


def draft_bullets(body_md: str) -> List[str]:
    match = re.search(
        r"\*\*右侧(?:\s*\d+)?\s*bullets\*\*(.*?)(?:\*\*讲者备注\*\*|$)",
        body_md,
        re.S | re.I,
    )
    if not match:
        return []
    rows = []
    for line in match.group(1).splitlines():
        stripped = line.strip()
        if stripped.startswith("-"):
            rows.append(_plain_text(stripped[1:]))
        elif re.match(r"^\d+\.\s+", stripped):
            rows.append(_plain_text(re.sub(r"^\d+\.\s+", "", stripped)))
    return [row for row in rows if row]


def draft_notes(body_md: str) -> str:
    parts: List[str] = []
    main_info = re.search(
        r"\*\*主信息\*\*\s*(.*?)(?:\*\*主图\*\*|\*\*右侧|\Z)",
        body_md,
        re.S,
    )
    if main_info and main_info.group(1).strip():
        parts.append(_plain_text(main_info.group(1)))
    notes = re.search(r"\*\*讲者备注\*\*(.*)$", body_md, re.S)
    if notes:
        for line in notes.group(1).splitlines():
            stripped = line.strip()
            if stripped.startswith("-"):
                parts.append(_plain_text(stripped[1:]))
    return "\n".join(part for part in parts if part)


def _close_open_target(path: Path) -> Tuple[bool, Optional[Path]]:
    if sys.platform != "win32":
        return False, None
    try:
        import win32com.client

        app = win32com.client.GetActiveObject("PowerPoint.Application")
        presentations = app.Presentations
        count = presentations.Count
    except Exception:
        return False, None

    target = os.path.normcase(os.path.abspath(path))
    for index in range(count, 0, -1):
        presentation = presentations.Item(index)
        if os.path.normcase(os.path.abspath(presentation.FullName)) != target:
            continue
        backup = None
        if presentation.Saved == 0:
            from datetime import datetime

            backup = path.with_stem(f"{path.stem}_prebuild_{datetime.now():%Y%m%d_%H%M%S}")
            presentation.SaveCopyAs(str(backup))
        presentation.Saved = True
        presentation.Close()
        return True, backup
    return False, None


def _open_output(path: Path) -> None:
    if sys.platform == "win32":
        os.startfile(str(path))


def _apply_sections(path: Path, specs: Sequence[Any], config: PptBuildConfig) -> None:
    if not config.sections or sys.platform != "win32":
        return
    slug_to_slide = {
        spec.slug: index + (2 if config.include_cover else 1)
        for index, spec in enumerate(specs)
    }
    sections = [
        (section.name, slug_to_slide[section.first_slug])
        for section in config.sections
        if section.first_slug in slug_to_slide
    ]
    if not sections:
        return
    if config.include_cover:
        sections[0] = (sections[0][0], 1)

    try:
        import win32com.client

        app = win32com.client.DispatchEx("PowerPoint.Application")
        app.Visible = True
        presentation = app.Presentations.Open(str(path), 0, 0, 0)
        properties = presentation.SectionProperties
        for index in range(properties.Count, 0, -1):
            try:
                properties.Delete(index, 0)
            except Exception:
                pass
        for name, start in sections:
            properties.AddBeforeSlide(start, name)
        presentation.Save()
        presentation.Close()
        app.Quit()
    except Exception as exc:
        print(f"WARN PowerPoint sections not applied: {exc}")


def _set_cover(slide, title: str, subtitle: str) -> None:
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_shape.text_frame.word_wrap = False
    for paragraph in title_shape.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def _add_header(slide, title: str, subtitle: str, config: PptBuildConfig):
    title_units = sum(1.0 if ord(char) > 127 else 0.55 for char in title)
    two_lines = title_units > 47
    title_y = Inches(0.03 if two_lines else 0.08)
    title_h = Inches(0.58 if two_lines else 0.34)
    subtitle_y = Inches(0.66 if two_lines else 0.47)
    divider_y = Inches(1.00 if two_lines else 0.81)
    body_top = Inches(1.15 if two_lines else 1.00)

    title_box = slide.shapes.add_textbox(
        Inches(0.40), title_y, SLIDE_W - Inches(0.80), title_h
    )
    title_box.text_frame.text = title
    title_box.text_frame.word_wrap = two_lines
    title_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    title_box.text_frame.margin_top = title_box.text_frame.margin_bottom = 0
    title_paragraph = title_box.text_frame.paragraphs[0]
    title_paragraph.alignment = PP_ALIGN.CENTER
    title_paragraph.line_spacing = 1.0
    for run in title_paragraph.runs:
        run.font.name = "微软雅黑"
        run.font.size = Pt(config.title_font_pt)
        run.font.bold = True

    subtitle_box = slide.shapes.add_textbox(
        Inches(0.06), subtitle_y, SLIDE_W - Inches(0.12), Inches(0.28)
    )
    subtitle_box.text_frame.text = subtitle
    subtitle_box.text_frame.word_wrap = False
    subtitle_box.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    subtitle_box.text_frame.margin_top = subtitle_box.text_frame.margin_bottom = 0
    subtitle_paragraph = subtitle_box.text_frame.paragraphs[0]
    subtitle_paragraph.alignment = PP_ALIGN.LEFT
    for run in subtitle_paragraph.runs:
        run.font.name = "微软雅黑"
        run.font.size = Pt(config.subtitle_font_pt)

    divider = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, divider_y, SLIDE_W, Inches(0.04)
    )
    divider.fill.solid()
    divider.fill.fore_color.rgb = DIVIDER_COLOR
    divider.line.fill.background()
    return body_top


def _bullet_lines(text: str, width_in: float, font_size: int) -> int:
    chars_per_line = max(2, int(width_in * 3.8 * 16 / font_size))
    return max(1, (len(text) + chars_per_line - 1) // chars_per_line)


def _bullet_height_in(items: Sequence[str], width_in: float, font_size: int) -> float:
    line_height = font_size / 72 * 1.45
    return 0.18 + sum(
        _bullet_lines(item, width_in, font_size) * line_height + 6 / 72
        for item in items
    )


def _add_bullets(slide, items: Sequence[str], x, y, width, height, font_size: int) -> None:
    if not items:
        return
    box = slide.shapes.add_textbox(x, y, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = "➤ " + item.strip()
        paragraph.space_before = Pt(6)
        paragraph.line_spacing = 1.5
        for run in paragraph.runs:
            run.font.name = "微软雅黑"
            run.font.size = Pt(font_size)


def _image_ratio(path: Path) -> float:
    with Image.open(path) as image:
        return image.width / image.height if image.height else 1.0


def _image_widths(images: Sequence[Path], total_width, explicit=None):
    if explicit and len(explicit) == len(images) and sum(explicit) > 0:
        weights = [float(value) / sum(explicit) for value in explicit]
    elif len(images) == 2:
        ratios = [_image_ratio(path) for path in images]
        first = min(0.62, max(0.38, ratios[0] / sum(ratios)))
        weights = [first, 1 - first]
    else:
        weights = [1 / len(images)] * len(images)
    return [total_width * weight for weight in weights]


def _add_image(slide, path: Path, x, y, max_width, max_height):
    ratio = _image_ratio(path)
    width = int(max_width)
    height = int(width / ratio)
    if height > int(max_height):
        height = int(max_height)
        width = int(height * ratio)
    left = int(x + (max_width - width) / 2)
    top = int(y + (max_height - height) / 2)
    return slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def _fitted_area(path: Path, width_in: float, height_in: float) -> float:
    ratio = _image_ratio(path)
    width = min(width_in, height_in * ratio)
    return width * (width / ratio)


def _score(
    image_area: float,
    text_fit: float,
    used_area: float,
    semantic: float,
    body_top=BODY_TOP,
) -> float:
    body_area = ((SLIDE_W - Inches(0.24)) / EMU_PER_IN) * (
        (BODY_BOTTOM - body_top) / EMU_PER_IN
    )
    image_readability = min(1.0, image_area / (body_area * 0.50))
    utilization = min(1.0, used_area / body_area)
    return (
        0.45 * image_readability
        + 0.25 * text_fit
        + 0.15 * utilization
        + 0.15 * semantic
        - 0.60 * max(0.0, 1.0 - text_fit)
    )


def _semantic_pair(images: Sequence[Path], bullets: Sequence[str]) -> bool:
    if len(images) != 2 or len(bullets) != 3:
        return False
    ignored = {
        "FIG", "FIGURE", "PLOT", "REPORT", "DOMAIN", "PEPTIDE",
        "PROTEOME", "METABOLOME", "QC",
    }
    tokens = []
    for path in images:
        tokens.append({
            token.upper()
            for token in re.findall(r"[A-Za-z0-9]+", path.stem)
            if len(token) >= 4 and token.upper() not in ignored
        })
    return all(any(token in bullets[index].upper() for token in tokens[index]) for index in range(2))


def _bottom_geometry(images, bullets, gap, explicit=None, body_top=BODY_TOP):
    font_size = 13
    note_width = (
        SLIDE_W - Inches(0.24) - gap * (len(bullets) - 1)
    ) / len(bullets)
    max_lines = max(
        _bullet_lines(item, note_width / EMU_PER_IN, font_size)
        for item in bullets
    )
    notes_height = Inches(
        min(2.15, max(1.20, 0.30 + max_lines * font_size / 72 * 1.45))
    )
    note_gap = Inches(0.10)
    total_width = SLIDE_W - Inches(0.24) - gap * (len(images) - 1)
    widths = _image_widths(images, total_width, explicit)
    display_height = min(
        BODY_BOTTOM - body_top - notes_height - note_gap,
        max(width / _image_ratio(path) for path, width in zip(images, widths)),
    )
    group_height = display_height + note_gap + notes_height
    group_top = body_top + (BODY_BOTTOM - body_top - group_height) / 2
    return widths, display_height, group_top, notes_height, note_gap


def _paired_geometry(images, bullets, gap, explicit=None, body_top=BODY_TOP):
    total_width = SLIDE_W - Inches(0.24) - gap
    widths = _image_widths(images, total_width, explicit)
    pair_lines = max(
        _bullet_lines(bullets[index], widths[index] / EMU_PER_IN, 13)
        for index in range(2)
    )
    pair_height = Inches(min(1.75, max(0.90, 0.28 + pair_lines * 13 / 72 * 1.45)))
    summary_lines = _bullet_lines(
        bullets[2], (SLIDE_W - Inches(0.24)) / EMU_PER_IN, 14
    )
    summary_height = Inches(
        min(1.05, max(0.55, 0.22 + summary_lines * 14 / 72 * 1.45))
    )
    inner_gap = Inches(0.08)
    note_gap = Inches(0.10)
    notes_height = pair_height + inner_gap + summary_height
    display_height = min(
        BODY_BOTTOM - body_top - notes_height - note_gap,
        max(width / _image_ratio(path) for path, width in zip(images, widths)),
    )
    group_height = display_height + note_gap + notes_height
    group_top = body_top + (BODY_BOTTOM - body_top - group_height) / 2
    return (
        widths, display_height, group_top, pair_height,
        summary_height, note_gap, inner_gap,
    )


def choose_layout(
    images: Sequence[Path],
    bullets: Sequence[str],
    *,
    dense: bool = False,
    explicit_widths=None,
    body_top=BODY_TOP,
) -> LayoutPlan:
    gap = Inches(0.12)
    body_height = (BODY_BOTTOM - body_top) / EMU_PER_IN
    right_width = (4.65 if dense else 3.12) if bullets else 0.0
    side_font = 13 if dense else 15
    side_total = (
        SLIDE_W
        - Inches(0.24 + right_width + (0.08 if bullets else 0.0))
        - gap * (len(images) - 1)
    )
    side_widths = _image_widths(images, side_total, explicit_widths)
    side_area = sum(
        _fitted_area(path, width / EMU_PER_IN, body_height)
        for path, width in zip(images, side_widths)
    )
    side_needed = _bullet_height_in(bullets, right_width, side_font)
    side_fit = min(1.0, body_height / side_needed)
    scores: Dict[str, float] = {
        "side": _score(
            side_area,
            side_fit,
            side_area + right_width * min(side_needed, body_height),
            0.60,
            body_top,
        )
    }

    # wide_bottom candidate: ratio>=2.5 (wide figures) OR side text overflows
    # (fit<1) — gives tall single-image pages with long bullets an escape
    # layout so text never renders past the divider/slide bottom.
    if len(images) == 1 and bullets and (
        _image_ratio(images[0]) >= 2.5 or side_fit < 1.0
    ):
        width = (SLIDE_W - Inches(0.56)) / EMU_PER_IN
        notes_needed = _bullet_height_in(bullets, width, 13)
        notes_height = min(2.10, max(0.90, notes_needed))
        image_height = min(4.95, body_height - notes_height - 0.12)
        image_area = _fitted_area(images[0], width, image_height)
        scores["wide_bottom"] = _score(
            image_area,
            min(1.0, notes_height / notes_needed),
            image_area + width * notes_height,
            0.80,
            body_top,
        )

    if len(images) == 2 and len(bullets) <= 3 and min(_image_ratio(p) for p in images) >= 1.30:
        widths, display_h, _, notes_h, _ = _bottom_geometry(
            images, bullets, gap, explicit_widths, body_top
        )
        image_area = sum(
            (width / EMU_PER_IN) * (display_h / EMU_PER_IN)
            for width in widths
        )
        note_width = (
            SLIDE_W - Inches(0.24) - gap * (len(bullets) - 1)
        ) / len(bullets)
        needed = max(
            _bullet_height_in([item], note_width / EMU_PER_IN, 13)
            for item in bullets
        )
        scores["bottom_columns"] = _score(
            image_area,
            min(1.0, (notes_h / EMU_PER_IN) / needed),
            image_area
            + ((SLIDE_W - Inches(0.24)) / EMU_PER_IN) * (notes_h / EMU_PER_IN),
            0.80,
            body_top,
        )

    if _semantic_pair(images, bullets):
        widths, display_h, _, pair_h, summary_h, _, _ = _paired_geometry(
            images, bullets, gap, explicit_widths, body_top
        )
        image_area = sum(
            (width / EMU_PER_IN) * (display_h / EMU_PER_IN)
            for width in widths
        )
        pair_needed = max(
            _bullet_height_in([bullets[index]], widths[index] / EMU_PER_IN, 13)
            for index in range(2)
        )
        summary_needed = _bullet_height_in(
            [bullets[2]], (SLIDE_W - Inches(0.24)) / EMU_PER_IN, 14
        )
        text_fit = min(
            1.0,
            (pair_h / EMU_PER_IN) / pair_needed,
            (summary_h / EMU_PER_IN) / summary_needed,
        )
        note_area = sum(
            (width / EMU_PER_IN) * (pair_h / EMU_PER_IN)
            for width in widths
        )
        note_area += (
            (SLIDE_W - Inches(0.24)) / EMU_PER_IN
        ) * (summary_h / EMU_PER_IN)
        scores["paired_summary"] = _score(
            image_area, text_fit, image_area + note_area, 1.0, body_top
        ) + 0.03

    name = max(scores, key=scores.get)
    return LayoutPlan(name=name, score=scores[name], scores=scores)


def _render_slide(slide, spec: Any, images: Sequence[Path], bullets: Sequence[str], config: PptBuildConfig) -> None:
    body_top = _add_header(slide, spec.title, spec.subtitle, config)
    if not images:
        _add_bullets(
            slide, bullets, Inches(0.95), body_top,
            Inches(11.40), BODY_BOTTOM - body_top, 16,
        )
        return

    dense = spec.slug in config.compact_slugs
    explicit = config.width_weights.get(spec.slug)
    plan = choose_layout(
        images,
        bullets,
        dense=dense,
        explicit_widths=explicit,
        body_top=body_top,
    )
    print(
        f"  layout {spec.title}: {plan.name} ("
        + ", ".join(f"{name}={score:.3f}" for name, score in plan.scores.items())
        + ")"
    )
    gap = Inches(0.12)
    left = Inches(0.12)

    if plan.name == "bottom_columns":
        widths, image_h, top, note_h, note_gap = _bottom_geometry(
            images, bullets, gap, explicit, body_top
        )
        x = left
        for path, width in zip(images, widths):
            _add_image(slide, path, x, top, width, image_h)
            x += width + gap
        note_width = (
            SLIDE_W - Inches(0.24) - gap * (len(bullets) - 1)
        ) / len(bullets)
        for index, bullet in enumerate(bullets):
            _add_bullets(
                slide, [bullet], left + index * (note_width + gap),
                top + image_h + note_gap, note_width, note_h, 13,
            )
        return

    if plan.name == "paired_summary":
        widths, image_h, top, pair_h, summary_h, note_gap, inner_gap = _paired_geometry(
            images, bullets, gap, explicit, body_top
        )
        positions = []
        x = left
        for path, width in zip(images, widths):
            positions.append(x)
            _add_image(slide, path, x, top, width, image_h)
            x += width + gap
        pair_y = top + image_h + note_gap
        for index in range(2):
            _add_bullets(
                slide, [bullets[index]], positions[index], pair_y,
                widths[index], pair_h, 13,
            )
        _add_bullets(
            slide, [bullets[2]], left, pair_y + pair_h + inner_gap,
            SLIDE_W - Inches(0.24), summary_h, 14,
        )
        return

    if plan.name == "wide_bottom":
        width = SLIDE_W - Inches(0.56)
        notes_needed = _bullet_height_in(bullets, width / EMU_PER_IN, 13)
        # tall-figure variant: cap notes at what fits after the image, never
        # let the note box start below BODY_BOTTOM - notes_h
        notes_h = Inches(min(2.40, max(0.90, notes_needed)))
        note_gap = Inches(0.12)
        image_h = min(Inches(4.95), BODY_BOTTOM - body_top - notes_h - note_gap)
        # if the estimated note height itself cannot fit in the body at all,
        # shrink font floor to 12pt before allowing overflow
        if image_h <= Inches(1.0) and notes_h < Inches(notes_needed if isinstance(notes_needed, float) else 0):
            notes_h = Inches(notes_needed)
        total_h = image_h + note_gap + notes_h
        top = body_top + (BODY_BOTTOM - body_top - total_h) / 2
        _add_image(slide, images[0], Inches(0.28), top, width, image_h)
        _add_bullets(
            slide, bullets, Inches(0.28), top + image_h + note_gap,
            width, notes_h, 13,
        )
        return

    right_width = Inches((4.65 if dense else 3.12) if bullets else 0.0)
    right_x = SLIDE_W - Inches(0.12) - right_width
    total_width = (
        right_x - left - (Inches(0.08) if bullets else 0) - gap * (len(images) - 1)
    )
    widths = _image_widths(images, total_width, explicit)
    x = left
    for path, width in zip(images, widths):
        _add_image(slide, path, x, body_top, width, BODY_BOTTOM - body_top)
        x += width + gap
    _add_bullets(
        slide, bullets, right_x, body_top, right_width,
        BODY_BOTTOM - body_top, 13 if dense else 15,
    )


def _resolve_images(spec: Any, config: PptBuildConfig) -> List[Path]:
    roots = [Path(config.figure_dir), *[Path(path) for path in config.figure_dirs]]
    images: List[Path] = []
    missing: List[str] = []
    for src in spec.images:
        name = Path(src).name
        resolved = next((root / name for root in roots if (root / name).is_file()), None)
        if resolved is None:
            missing.append(name)
        else:
            images.append(resolved)
    if missing:
        raise FileNotFoundError(
            f"{spec.slug}: missing PPT figures {missing}; searched {[str(root) for root in roots]}"
        )
    return images


def build_draft_ppt(specs: Iterable[Any], config: PptBuildConfig) -> Path:
    specs = list(specs)
    template = _resolve_template(config.template_candidates)
    prs = _load_template(template)
    blank_layout = prs.slide_layouts[6]

    if config.include_cover:
        cover = prs.slides.add_slide(prs.slide_layouts[0])
        _set_cover(cover, config.title, config.subtitle)

    for spec in specs:
        slide = prs.slides.add_slide(blank_layout)
        images = _resolve_images(spec, config)
        bullets = draft_bullets(spec.body_md)
        _render_slide(slide, spec, images, bullets, config)
        notes = draft_notes(spec.body_md)
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    output = Path(config.output)
    output.parent.mkdir(parents=True, exist_ok=True)
    temporary = output.with_name(f".{output.stem}.build_tmp{output.suffix}")
    temporary.unlink(missing_ok=True)
    prs.save(str(temporary))

    closed, backup = _close_open_target(output)
    os.replace(temporary, output)
    _apply_sections(output, specs, config)
    if config.auto_open:
        _open_output(output)

    if closed:
        print(f"Closed and reopened target PPT: {output}")
    if backup:
        print(f"Backed up unsaved target PPT: {backup}")
    print(f"Wrote draft PPT: {output}")
    print(f"  slides: {len(prs.slides)}")
    return output
