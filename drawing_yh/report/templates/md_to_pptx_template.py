#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
md_to_pptx_template.py — REPORT.md → native PPTX,继承一份 PPTX 母版的 theme。

复制到项目目录改名 `md_to_pptx.py`。配置(SLIDES)从 `_report_config.py` import。
**改这 4 处即可**:
  1. PROJECT_LABEL —— 顶部 header 左上的小字(类似 "Retinal (Macaca, bulk)")
  2. TEMPLATE      —— 母版 .pptx 路径(用它的 master / 字体 / 主题色;原文件不动)
  3. OUT           —— 输出 .pptx 路径
  4. main() 里的 title_slide 标题 + 副标题文案

布局:每张 slide 顶部 = 项目标签 + slide 标题 + 橙色分隔线;body 自动选择
  - 有图(`<img>` 或 `![caption](path)`):左大图 + 右 ➤ bullets(垂直居中)
  - 多图:默认并排 / 2×2 排布,也可用 `_report_config.PPT_SLIDES` 指定
  - 无图:全宽 ➤ bullets / 表
约定:有图 slide 不放表;表格列宽按内容长短分配;字号 16pt(表 14pt header / 11–14pt body);
      内容超 → 自动截字号 / 截行 / 丢次要 list,**不**显式提示。

注:本模板假设母版尺寸 16:9(13.33 × 7.5 in)。其他比例改下方 SLIDE_W / SLIDE_H。
"""
import re, sys
from pathlib import Path
from urllib.parse import unquote, urlparse
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ----- 改这里 -----
PROJECT_LABEL  = 'My Project × Subtopic'                   # 改成你项目的标签
DIVIDER_COLOR  = RGBColor(0xBE, 0x76, 0x00)                # header 下方分隔线颜色
TEMPLATE = Path(r'D:/path/to/your/template.pptx')          # 母版 PPTX
OUT      = Path(r'D:/path/to/output/REPORT.pptx')          # 输出
# -------------------

# 共享配置(纯 import 不触发任何渲染副作用)
sys.path.insert(0, str(Path(__file__).resolve().parent))
from _report_config import HERE, SLIDES, load_chunks   # noqa: E402
try:
    from _report_config import FIGURE_SRC              # noqa: E402
except ImportError:
    FIGURE_SRC = None
try:
    from _report_config import PPT_IMAGE_DIRS          # noqa: E402
except ImportError:
    PPT_IMAGE_DIRS = []
try:
    from _report_config import PPT_SLIDES              # noqa: E402
except ImportError:
    PPT_SLIDES = {}

# 母版 16:9(13.33 × 7.5 in),hardcode 避免 module-level 创建 Presentation
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.500)

# prs / 各 layout 在 main() 里加载,通过 _PRS_CTX 给 add_content_slide 用
_PRS_CTX: dict = {}     # {'blank_layout': layout 对象}
_PPT_CTX: dict = {}     # 当前 slide 上下文,用于诊断
_PPT_DIAG: dict[str, list[str]] = {
    'missing_chunks': [],
    'missing_images': [],
    'unsupported_images': [],
    'fallback_images': [],
    'truncated_images': [],
    'dropped_tables': [],
}

SUPPORTED_IMAGE_EXTS = {'.png', '.jpg', '.jpeg', '.bmp', '.gif', '.tif', '.tiff'}
RASTER_FALLBACK_EXTS = ('.png', '.jpg', '.jpeg')


def _diag(kind: str, msg: str):
    bucket = _PPT_DIAG.setdefault(kind, [])
    if msg not in bucket:
        bucket.append(msg)


def _print_diagnostics():
    any_warn = False
    for kind, rows in _PPT_DIAG.items():
        if not rows:
            continue
        any_warn = True
        print(f'PPT diagnostics [{kind}]:')
        for row in rows:
            print(f'  - {row}')
    if not any_warn:
        print('PPT diagnostics: no warnings')


def _load_template() -> Presentation:
    """读母版 + 真删原 slide(连 part + rels + 图片一起删,否则 PowerPoint
    严格校验会报"内容有问题")。继承 master / theme / 字体配色。"""
    prs = Presentation(str(TEMPLATE))
    sld_id_lst = prs.slides._sldIdLst
    for sld_id in list(sld_id_lst):
        prs.part.drop_rel(sld_id.rId)
        sld_id_lst.remove(sld_id)
    return prs


_HTML_IMG_RE = re.compile(r'<img\b[^>]*\bsrc\s*=\s*["\']([^"\']+)["\']', re.I)
_MD_IMG_RE = re.compile(r'!\[[^\]]*\]\(([^)]+)\)')


def _extract_md_image_src(part: str) -> str | None:
    m = _MD_IMG_RE.search(part)
    if not m:
        return None
    raw = m.group(1).strip()
    if raw.startswith('<') and raw.endswith('>'):
        raw = raw[1:-1].strip()
    # 去掉可选 title: ![](fig.png "caption")
    m_title = re.match(r'([^\s]+)\s+["\'][^"\']*["\']$', raw)
    if m_title:
        raw = m_title.group(1)
    return raw or None


# ------------------------------------------------------------------
# 2. Markdown chunk → 块结构(para / list / table / image / code)
# ------------------------------------------------------------------
def parse_chunk(md: str):
    """返回 (heading, [block, ...])。heading 用 chunk 第一个 H2/H3。"""
    h_m = re.search(r'^(#{2,3})\s+(.+?)\s*$', md, re.MULTILINE)
    heading = h_m.group(2).strip() if h_m else ''
    body = (md[h_m.end():] if h_m else md).strip()

    blocks = []
    for part in re.split(r'\n\s*\n', body):
        part = part.strip()
        if not part:
            continue
        # 跳过 markdown 横线(---、***、___)
        if re.fullmatch(r'[-*_]{3,}', part):
            continue
        # 跳过 image 后面的 *italic caption*(我们的左大图布局不需要再放图注)
        if re.fullmatch(r'\*[^*]+\*', part):
            continue
        if part.startswith('|') and re.search(r'\n\|', part):
            blocks.append({'type': 'table', 'md': part})
        elif re.match(r'<img\s', part, re.I):
            m = _HTML_IMG_RE.search(part)
            if m:
                blocks.append({'type': 'image', 'src': m.group(1)})
        elif re.match(r'!\[', part):
            src = _extract_md_image_src(part)
            if src:
                blocks.append({'type': 'image', 'src': src})
        elif part.startswith('```'):
            inner = re.sub(r'^```\w*\n?', '', part)
            inner = re.sub(r'\n?```$', '', inner)
            blocks.append({'type': 'code', 'text': inner})
        elif re.match(r'[-*]\s', part):
            items = [re.sub(r'^[-*]\s+', '', l).strip()
                     for l in part.splitlines() if l.strip().startswith(('-', '*'))]
            blocks.append({'type': 'list', 'items': items})
        else:
            blocks.append({'type': 'para', 'text': part})
    return heading, blocks


# ------------------------------------------------------------------
# 3. Helper:在 slide 上的 (cur_y) 位置追加各种块
# ------------------------------------------------------------------
def _iter_image_dirs():
    dirs = [HERE / 'report_figs', HERE]
    if FIGURE_SRC:
        dirs.append(Path(FIGURE_SRC))
    for d in PPT_IMAGE_DIRS or []:
        dirs.append(Path(d))
    seen = set()
    for d in dirs:
        try:
            key = str(d.resolve())
        except OSError:
            key = str(d)
        if key in seen:
            continue
        seen.add(key)
        yield d


def _clean_src(src: str) -> str:
    src = unquote(str(src).strip())
    src = src.split('#', 1)[0].split('?', 1)[0]
    if src.startswith('<') and src.endswith('>'):
        src = src[1:-1].strip()
    return src


def _candidate_paths(src: str, *, raster_only: bool):
    clean = _clean_src(src)
    if not clean:
        return []
    parsed = urlparse(clean)
    if parsed.scheme in {'http', 'https'}:
        return []

    p = Path(clean)
    names = [p.name]
    if raster_only:
        names = [p.with_suffix(ext).name for ext in RASTER_FALLBACK_EXTS]

    candidates = []
    if p.is_absolute():
        candidates.append(p)
        if raster_only:
            candidates.extend(p.with_suffix(ext) for ext in RASTER_FALLBACK_EXTS)
    else:
        candidates.append((HERE / p).resolve())
        for d in _iter_image_dirs():
            for name in names:
                candidates.append(d / name)

    out = []
    seen = set()
    for c in candidates:
        s = str(c)
        if s not in seen:
            out.append(c)
            seen.add(s)
    return out


def resolve_image(src: str) -> Path | None:
    """解析本地 raster 图片。SVG/PDF 等非 raster 需要同名 PNG/JPG fallback。"""
    slug = _PPT_CTX.get('slug', '?')
    clean = _clean_src(src)
    parsed = urlparse(clean)
    if parsed.scheme in {'http', 'https'}:
        _diag('unsupported_images', f'{slug}: remote image is not embedded in PPTX: {src}')
        return None

    exact_existing = []
    for c in _candidate_paths(clean, raster_only=False):
        if c.exists():
            exact_existing.append(c)
            if c.suffix.lower() in SUPPORTED_IMAGE_EXTS:
                return c

    for c in _candidate_paths(clean, raster_only=True):
        if c.exists() and c.suffix.lower() in SUPPORTED_IMAGE_EXTS:
            if exact_existing:
                _diag('fallback_images', f'{slug}: used raster fallback {c.name} for {Path(clean).name}')
            return c

    if exact_existing:
        names = ', '.join(p.name for p in exact_existing)
        _diag('unsupported_images', f'{slug}: {names} found but PPTX needs PNG/JPG/BMP/GIF/TIFF')
    else:
        _diag('missing_images', f'{slug}: {src}')
    return None


BODY_LEFT   = Inches(0.30)
BODY_RIGHT  = Inches(0.30)
BODY_TOP    = Inches(1.00)              # 在 divider(y=0.81)之下留 gap
BODY_W      = SLIDE_W - BODY_LEFT - BODY_RIGHT
BODY_BOTTOM = SLIDE_H - Inches(0.25)


# ------------------------------------------------------------------
# Header layout 仿 retina-bulk:左上项目标签 + 中右 slide 标题 + 横分隔线
# ------------------------------------------------------------------
def add_header(slide, slide_title: str):
    # 左上 项目标签
    pl = slide.shapes.add_textbox(Inches(0.06), Inches(0.41), Inches(5.0), Inches(0.40))
    pl.text_frame.text = PROJECT_LABEL
    pl.text_frame.margin_top = pl.text_frame.margin_bottom = 0
    p = pl.text_frame.paragraphs[0]
    for r in p.runs:
        r.font.name = '微软雅黑'
        r.font.size = Pt(18)
    # 右半 slide 标题(bigger + bold)
    st = slide.shapes.add_textbox(Inches(5.20), Inches(0.13), Inches(8.00), Inches(0.55))
    st.text_frame.text = slide_title
    st.text_frame.margin_top = st.text_frame.margin_bottom = 0
    p = st.text_frame.paragraphs[0]
    for r in p.runs:
        r.font.name = '微软雅黑'
        r.font.size = Pt(28)
        r.font.bold = True
    # 横分隔线 #BE7600,full width × 0.04 in
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.81), SLIDE_W, Inches(0.04))
    div.fill.solid()
    div.fill.fore_color.rgb = DIVIDER_COLOR
    div.line.fill.background()  # 无边框


def add_paragraph(slide, txt: str, x, y, w, *, font_size=16, bold=False):
    # 16pt 中文一行约 25–30 字宽 5 in;每行 ~0.32 in 高
    chars_per_line = max(15, int(w / Emu(914400) * 5))   # rough
    n_lines = max(1, len(txt) // chars_per_line + 1)
    h = Inches(max(0.4, 0.34 * n_lines + 0.18))
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = txt
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = '微软雅黑'
            r.font.size = Pt(font_size)
            r.font.bold = bold
    return h


def add_list(slide, items, x, y, w, *, font_size=14):
    n = len(items)
    h = Inches(max(0.45, 0.32 * n + 0.2))
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = '• ' + it
        for r in p.runs:
            r.font.name = '微软雅黑'
            r.font.size = Pt(font_size)
    return h


def _bullet_lines(text: str, w_in: float, font_size: int) -> int:
    """估算文字在宽度 w_in (英寸) 的 textbox 里能占多少行。
    标定:微软雅黑 16pt,5 in 宽 ≈ 19 中文字/行 → 3.8 字/(in × 16pt)
    其他字号线性外推:chars_per_line = w_in × 3.8 × 16/font_size。"""
    chars_per_line = max(2, int(w_in * 3.8 * 16 / font_size))
    return max(1, (len(text) + chars_per_line - 1) // chars_per_line)


def _parse_md_table(md_block: str):
    """返回 (cells, n_cols)。cells 是 list of list of str。"""
    rows = [r for r in md_block.splitlines() if r.strip().startswith('|')
            and not re.match(r'^\|\s*[-:|\s]+\|\s*$', r)]
    cells = [[c.strip() for c in r.strip().strip('|').split('|')] for r in rows]
    n_cols = max((len(r) for r in cells), default=0)
    return cells, n_cols


def _compute_col_widths_in(cells, total_w_in: float, n_cols: int,
                            min_col_w_in: float = 0.6) -> list[float]:
    """按每列 max(len(cell)) 加权分配列宽,带 min 兜底。
    充分利用版面 → 短列窄、长列宽,减少不必要的 wrap。"""
    if n_cols == 0: return []
    max_chars = [1] * n_cols
    for row in cells:
        for ci in range(min(len(row), n_cols)):
            max_chars[ci] = max(max_chars[ci], len(row[ci]))
    # min 兜底,但不能加起来超过总宽
    min_alloc = min(min_col_w_in, total_w_in / n_cols)
    total_min = min_alloc * n_cols
    if total_min >= total_w_in:
        return [total_w_in / n_cols] * n_cols
    remaining = total_w_in - total_min
    total_wt = sum(max_chars)
    return [min_alloc + (max_chars[i] / total_wt) * remaining for i in range(n_cols)]


def _cells_h_in(cells, col_widths_in: list[float], font_pt: int) -> float:
    """按每列实际宽度算 wrap 高度(英寸)。"""
    if not cells: return 0.0
    line_h_in = (font_pt + 6) / 72.0
    total_h = 0.05
    for row in cells:
        max_lines = 1
        for ci in range(min(len(row), len(col_widths_in))):
            cw_in = col_widths_in[ci]
            chars_per_line = max(2, int(cw_in * 3.8 * 16 / font_pt))
            n = max(1, (len(row[ci]) + chars_per_line - 1) // chars_per_line)
            max_lines = max(max_lines, n)
        total_h += max_lines * line_h_in + 0.06
    return total_h


def estimate_table_h_in(md_block: str, w_in: float, font_pt: int = 14) -> float:
    cells, n_cols = _parse_md_table(md_block)
    if not cells: return 0.0
    col_widths = _compute_col_widths_in(cells, w_in, n_cols)
    return _cells_h_in(cells, col_widths, font_pt)


def can_fit_table(md_block: str, w_in: float, max_h_in: float,
                   min_font_pt: int = 11) -> bool:
    """min_font 也装不下 → False。"""
    return estimate_table_h_in(md_block, w_in, font_pt=min_font_pt) <= max_h_in


def fit_bullets(items, w_in: float, max_h_in: float, font_size: int = 16):
    """在 (w × max_h) 的盒子里能塞下几条箭头要点 + 每条多长。
    超出的:先尝试缩字数,实在装不下就丢弃。返回 (items_fit, est_h_in)。"""
    line_h = font_size / 72 * 1.45                  # 行高 in
    space_before = 6 / 72                            # 段前空 in
    used = 0.18                                      # 上下 padding
    out = []
    for it in items:
        s = it.strip()
        n_lines = _bullet_lines(s, w_in, font_size)
        item_h = line_h * n_lines + space_before
        if used + item_h > max_h_in:
            # 试试压短到能放进去
            rem_h = max_h_in - used - space_before
            if rem_h < line_h:
                break                                 # 没空间,直接停
            allowed_lines = max(1, int(rem_h / line_h))
            chars_per_line = max(2, int(w_in * 3.8 * 16 / font_size))
            allowed_chars = allowed_lines * chars_per_line - 2   # 留省略号
            if allowed_chars >= 8:
                out.append(s[:allowed_chars].rstrip() + '…')
                used += line_h * allowed_lines + space_before
            break
        out.append(s)
        used += item_h
    return out, used


def add_arrow_bullets_centered(slide, items, x, y, w, h, *, font_size=16):
    """➤ bullets,在 (w × h) 盒子内**垂直居中**(用 PPT 的 vertical_anchor)。"""
    if not items:
        return Inches(0)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = '➤ ' + it.strip()
        p.space_before = Pt(6)
        for r in p.runs:
            r.font.name = '微软雅黑'
            r.font.size = Pt(font_size)
    return h


def add_arrow_bullets(slide, items, x, y, w, *, font_size=16, max_h=None):
    """➤ 箭头要点。可选 max_h 自动裁切(防溢出)。"""
    w_in = w / 914400  # EMU → in
    if max_h is not None:
        items, est_used = fit_bullets(items, w_in, max_h / 914400, font_size=font_size)
        h = Inches(min(est_used + 0.05, max_h / 914400))
    else:
        # 无限制 → 老式估高
        h = Inches(max(0.45, 0.42 * len(items) + 0.20))
    if not items:
        return Inches(0)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = '➤ ' + it.strip()
        p.space_before = Pt(6)
        for r in p.runs:
            r.font.name = '微软雅黑'
            r.font.size = Pt(font_size)
    return h


def add_md_table(slide, md_block: str, x, y, w, *, max_h=None, font_pt=14, min_font_pt=11):
    """画表,wrap-aware。
    流程:
      ① 在 [min_font, font_pt] 找最大可读字号能装下全表(含 wrap)
      ② 装不下 → 用 min_font 一行一行加,直到再加就溢出 → 末行 "(共 N 行,省略 M)"
      ③ header 至少要装下,否则返回 0(放弃这张表)
    """
    cells, n_cols = _parse_md_table(md_block)
    total_rows = len(cells)
    if not (total_rows and n_cols):
        return Inches(0)

    w_in = w / 914400
    # 关键:列宽按内容长短分配(短列窄、长列宽,充分利用版面)
    col_widths_in = _compute_col_widths_in(cells, w_in, n_cols)

    chosen_font = font_pt
    final_cells = cells

    if max_h is not None:
        max_h_in = max_h / 914400

        # ① 在大 → 小字号里找能装下全表的
        fits_full = False
        for fp in range(font_pt, min_font_pt - 1, -1):
            if _cells_h_in(cells, col_widths_in, fp) <= max_h_in:
                chosen_font = fp
                fits_full = True
                break

        # ② 全表装不下 → min font 截行
        if not fits_full:
            chosen_font = min_font_pt
            kept = [cells[0]]                       # 必带 header
            ellipsis_h = (chosen_font + 6) / 72.0 + 0.06
            for row in cells[1:]:
                trial = kept + [row]
                trial_h = _cells_h_in(trial, col_widths_in, chosen_font)
                if trial_h + ellipsis_h > max_h_in:
                    break
                kept = trial
            n_kept_data = len(kept) - 1
            n_dropped = (total_rows - 1) - n_kept_data
            if n_dropped > 0:
                ell = ['…'] * n_cols
                ell[1 if n_cols > 1 else 0] = f'(共 {total_rows - 1} 行,省略 {n_dropped})'
                kept.append(ell[:n_cols])
            # ③ 还是装不下(光 header 都超):整个放弃
            if _cells_h_in(kept, col_widths_in, chosen_font) > max_h_in:
                _diag('dropped_tables', f"{_PPT_CTX.get('slug', '?')}: table dropped; even header does not fit")
                return Inches(0)
            final_cells = kept
            n_kept_data = len(final_cells) - 2 if final_cells and final_cells[-1][0] == '…' else len(final_cells) - 1
            n_dropped = max(0, total_rows - 1 - n_kept_data)
            if n_dropped:
                _diag('dropped_tables', f"{_PPT_CTX.get('slug', '?')}: table truncated, omitted {n_dropped} rows")

    n_rows = len(final_cells)
    h = Inches(_cells_h_in(final_cells, col_widths_in, chosen_font))

    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, w, h).table
    # 应用每列宽度(short cell narrow,long cell wide)
    for ci, cw_in in enumerate(col_widths_in):
        if ci < n_cols:
            tbl.columns[ci].width = Inches(cw_in)
    for ri, row in enumerate(final_cells):
        for ci in range(n_cols):
            val = row[ci] if ci < len(row) else ''
            cell = tbl.cell(ri, ci)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = '微软雅黑'
                    r.font.size = Pt(chosen_font + 1 if ri == 0 else chosen_font)
                    if ri == 0:
                        r.font.bold = True
    return h


def add_image(slide, img_path: Path, x, y, max_w, max_h):
    """嵌图,先用 Pillow 读真实像素,等比 fit 到 (max_w × max_h),
    再在盒子内居中;一次 add_picture 给定确切宽高 —— 不留孤儿 relationship。"""
    with Image.open(img_path) as im:
        iw, ih = im.size
    aspect = ih / iw if iw else 1
    tw = int(max_w)
    th = int(tw * aspect)
    if th > int(max_h):
        th = int(max_h)
        tw = int(th / aspect)
    cx = int(x + (max_w - tw) / 2)
    cy = int(y + (max_h - th) / 2)
    pic = slide.shapes.add_picture(str(img_path), cx, cy, width=tw, height=th)
    return pic.height


def add_code(slide, txt: str, x, y, w):
    n_lines = txt.count('\n') + 1
    h = Inches(max(0.4, 0.22 * n_lines + 0.2))
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.text = txt
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.name = 'Consolas'
            r.font.size = Pt(11)
    return h


# ------------------------------------------------------------------
# 4. 一张 slide 的完整布局
# ------------------------------------------------------------------
# ------------------------------------------------------------------
# 双栏布局参数(有图时用):左大图 + 右箭头要点
# ------------------------------------------------------------------
LEFT_IMG_X   = Inches(0.30)
LEFT_IMG_W   = Inches(7.50)         # ~56% slide width
RIGHT_TXT_X  = Inches(7.95)
RIGHT_TXT_W  = Inches(5.10)         # ~38% slide width


def _collect_text_blocks(blocks):
    bullets, tables, codes = [], [], []
    for blk in blocks:
        if blk['type'] == 'para':
            bullets.append(blk['text'])
        elif blk['type'] == 'list':
            bullets.extend(blk['items'])
        elif blk['type'] == 'table':
            tables.append(blk)
        elif blk['type'] == 'code':
            codes.append(blk)
    return bullets, tables, codes


def layout_image_left(slide, images, others, *, bullets_override=None, text_heavy=False):
    """有图 → 左 7.5 in 大图,右 5.1 in 箭头要点 + 表格。"""
    left_img_w = Inches(7.20) if text_heavy else LEFT_IMG_W
    right_txt_x = Inches(7.55) if text_heavy else RIGHT_TXT_X
    right_txt_w = Inches(5.35) if text_heavy else RIGHT_TXT_W
    bullet_font_size = 14 if text_heavy else 16

    # 1. 左主图:取第一张可解析图片,非 raster 自动尝试同名 PNG/JPG fallback
    main_img = None
    for im in images:
        main_img = resolve_image(im['src'])
        if main_img:
            break
    if main_img:
        avail_h = BODY_BOTTOM - BODY_TOP
        add_image(slide, main_img, LEFT_IMG_X, BODY_TOP, left_img_w, avail_h)

    # 2. 右栏:把 para + list 全收编成 ➤ bullets,table 紧跟其后
    if bullets_override is not None:
        bullets = list(bullets_override)
        tables, codes = [], []
    else:
        bullets, tables, codes = _collect_text_blocks(others)

    # 右栏内容:**有图的 slide 一律不放表**(易冲突 / 拥挤)。
    # 表里的数据用 bullets 凝练,或读者去网页 / REPORT.md 看
    avail = BODY_BOTTOM - BODY_TOP                       # ~6.25 in
    avail_in = avail / 914400
    rt_w_in  = right_txt_w / 914400

    fit_table = None
    dropped_n = len(tables)

    # 2a) 只有 bullets,没表 → 垂直居中
    if bullets and fit_table is None:
        # 留底部 ~0.3 in 给"省略"注脚
        anchor_h = avail - (Inches(0.35) if dropped_n else Inches(0))
        # 截断超长 bullet
        items_fit, _ = fit_bullets(bullets, rt_w_in, anchor_h / 914400, font_size=bullet_font_size)
        add_arrow_bullets_centered(slide, items_fit, right_txt_x, BODY_TOP,
                                    right_txt_w, anchor_h, font_size=bullet_font_size)
        cur_y = BODY_BOTTOM - (Inches(0.35) if dropped_n else Inches(0))

    # 2b) bullets + 1 张表 → bullets 顶上,表跟下面
    elif fit_table is not None:
        cur_y = BODY_TOP
        if bullets:
            cur_y += add_arrow_bullets(slide, bullets, right_txt_x, cur_y, right_txt_w,
                                        font_size=bullet_font_size, max_h=avail * 0.35) + Inches(0.1)
        rem = BODY_BOTTOM - cur_y - Inches(0.4 if dropped_n else 0.05)
        cur_y += add_md_table(slide, fit_table['md'], right_txt_x, cur_y, right_txt_w,
                               max_h=rem, font_pt=14, min_font_pt=11) + Inches(0.05)

    # 2c) 没 bullets 也没能放的表
    else:
        cur_y = BODY_BOTTOM

    # 4) 代码块(如果还有空间,塞右下)
    for cb in codes:
        if cur_y >= BODY_BOTTOM - Inches(0.4): break
        cur_y += add_code(slide, cb['text'], right_txt_x, cur_y, right_txt_w) + Inches(0.05)

    dropped = max(0, len(images) - 1)
    if dropped:
        _diag('truncated_images', f"{_PPT_CTX.get('slug', '?')}: {dropped} extra image(s) ignored by image_left layout")


def add_multi_image_slide(prs, label: str, img_paths: list[Path], *, bullets=None,
                          width_weights=None, layout='auto'):
    """多图页:2 张横排;3-4 张走 2×2;超过 4 张只放前 4 张并诊断。"""
    slide = prs.slides.add_slide(_PRS_CTX['blank_layout'])
    add_header(slide, label)
    imgs = [p for p in img_paths if p]
    if not imgs:
        return
    if len(imgs) > 4:
        _diag('truncated_images', f"{_PPT_CTX.get('slug', '?')}: used first 4 of {len(imgs)} images")
        imgs = imgs[:4]

    right_w = Inches(2.35) if bullets else Inches(0)
    right_x = SLIDE_W - BODY_RIGHT - right_w
    left_x = BODY_LEFT
    top_y = BODY_TOP
    avail_w = right_x - left_x - (Inches(0.12) if bullets else Inches(0))
    avail_h = BODY_BOTTOM - BODY_TOP
    gap = Inches(0.16)
    n = len(imgs)

    if layout == 'grid' or n in (3, 4):
        cols = 2
        rows = 2 if n > 2 else 1
        cell_w = (avail_w - gap * (cols - 1)) / cols
        cell_h = (avail_h - gap * (rows - 1)) / rows
        for idx, ip in enumerate(imgs):
            row, col = divmod(idx, cols)
            add_image(slide, ip, left_x + col * (cell_w + gap), top_y + row * (cell_h + gap),
                      cell_w, cell_h)
    else:
        total_w = avail_w - gap * (n - 1)
        if width_weights and len(width_weights) == n and sum(width_weights) > 0:
            weights = [float(w) / sum(width_weights) for w in width_weights]
            widths = [total_w * w for w in weights]
        else:
            widths = [total_w / n] * n
        x = left_x
        for ip, w in zip(imgs, widths):
            add_image(slide, ip, x, top_y, w, avail_h)
            x += w + gap

    if bullets:
        add_arrow_bullets_centered(slide, bullets, right_x, BODY_TOP, right_w,
                                   BODY_BOTTOM - BODY_TOP, font_size=15)


def _estimate_height_in(blk, w_in: float = 12.73) -> float:
    """估高(英寸),用 wrap-aware 算法预判会不会溢出。"""
    if blk['type'] == 'list':
        return max(0.45, 0.40 * len(blk['items']) + 0.25)
    if blk['type'] == 'para':
        return max(0.4, 0.04 * (len(blk['text']) // 70 + 1) + 0.32)
    if blk['type'] == 'table':
        return estimate_table_h_in(blk['md'], w_in, font_pt=14)
    if blk['type'] == 'code':
        return max(0.4, 0.22 * (blk['text'].count('\n') + 1) + 0.2)
    return 0.5


def layout_full_width(slide, blocks):
    """无图 → body 全宽。**有表的 slide 表为主**:
       - bullet list 整个丢(那些"观察 / 评论"放网页/md 看)
       - paragraph 截到 60 字(基本 1 行作 label)
       - 表用 wrap-aware 自动选字号,装不下截行
    """
    n_tables = sum(1 for b in blocks if b['type'] == 'table')
    table_dominated = n_tables > 0
    para_max_chars = 60 if table_dominated else 1000
    body_w_in = BODY_W / 914400

    # 表为主时 → 直接过滤 list 块(整段丢)
    n_dropped_lists = 0
    if table_dominated:
        n_dropped_lists = sum(1 for b in blocks if b['type'] == 'list')
        blocks = [b for b in blocks if b['type'] != 'list']

    seen_tbl = 0
    dropped_tbl = 0
    cur_y = BODY_TOP

    for blk in blocks:
        rem = BODY_BOTTOM - cur_y - Inches(0.05)
        if rem < Inches(0.4):
            if blk['type'] == 'table': dropped_tbl += 1
            continue
        if blk['type'] == 'list':
            cur_y += add_arrow_bullets(slide, blk['items'], BODY_LEFT, cur_y, BODY_W,
                                        font_size=16, max_h=rem) + Inches(0.06)
        elif blk['type'] == 'para':
            text = blk['text']
            if len(text) > para_max_chars:
                text = text[:para_max_chars - 1].rstrip() + '…'
            items, _ = fit_bullets([text], body_w_in, rem / 914400, font_size=16)
            if items:
                cur_y += add_arrow_bullets(slide, items, BODY_LEFT, cur_y, BODY_W,
                                            font_size=16, max_h=rem) + Inches(0.04)
        elif blk['type'] == 'table':
            if seen_tbl >= 2:
                _diag('dropped_tables', f"{_PPT_CTX.get('slug', '?')}: extra table dropped")
                dropped_tbl += 1
                continue
            h_used = add_md_table(slide, blk['md'], BODY_LEFT, cur_y, BODY_W,
                                   max_h=rem, font_pt=14, min_font_pt=10)
            if h_used == 0:
                dropped_tbl += 1
            else:
                cur_y += h_used + Inches(0.08)
                seen_tbl += 1
        elif blk['type'] == 'code':
            cur_y += add_code(slide, blk['text'], BODY_LEFT, cur_y, BODY_W) + Inches(0.05)



def _images_of_chunk(md_chunk: str) -> list[dict]:
    _, blocks = parse_chunk(md_chunk)
    return [b for b in blocks if b['type'] == 'image']


def _resolved_images(images: list[dict]) -> list[Path]:
    out = []
    for b in images:
        ip = resolve_image(b['src'])
        if ip:
            out.append(ip)
    return out


def _normalise_slide_opts(opts):
    if opts is None:
        return {}
    if isinstance(opts, str):
        return {'layout': opts}
    return dict(opts)


def _slug_slide_spec(slug: str, label: str, cat_key: str):
    opts = _normalise_slide_opts(PPT_SLIDES.get(slug) if isinstance(PPT_SLIDES, dict) else None)
    return {
        'slug': slug,
        'chunk': slug,
        'label': opts.pop('label', label),
        'category': cat_key,
        **opts,
    }


def _iter_ppt_specs():
    """默认按 SLIDES 输出;若 PPT_SLIDES 是 list,则完全按该 list 选页和排序。"""
    if isinstance(PPT_SLIDES, list):
        for item in PPT_SLIDES:
            if isinstance(item, str):
                # 用 SLIDES 里的 label/category 补齐
                for slug, _, label, cat_key in SLIDES:
                    if slug == item:
                        yield _slug_slide_spec(slug, label, cat_key)
                        break
                else:
                    yield {'slug': item, 'chunk': item, 'label': item, 'category': None}
            elif isinstance(item, dict):
                spec = dict(item)
                spec.setdefault('slug', spec.get('chunk'))
                spec.setdefault('chunk', spec['slug'])
                spec.setdefault('label', spec['slug'])
                yield spec
        return

    for slug, _, label, cat_key in SLIDES:
        yield _slug_slide_spec(slug, label, cat_key)


def _select_named_images(images: list[dict], names) -> list[dict]:
    if not names:
        return images
    wanted = {Path(str(n)).stem for n in names}
    out = []
    for b in images:
        stem = Path(_clean_src(b['src'])).stem
        if stem in wanted:
            out.append(b)
    return out


def add_content_slide(prs, label: str, md_chunk: str, *, spec=None):
    spec = spec or {}
    _, blocks = parse_chunk(md_chunk)
    images = [b for b in blocks if b['type'] == 'image']
    others = [b for b in blocks if b['type'] != 'image']
    images = _select_named_images(images, spec.get('images'))
    bullets_override = spec.get('bullets')
    text_heavy = bool(spec.get('text_heavy', False))
    layout = spec.get('layout', 'auto')

    if images and (
        layout in {'multi_images', 'two_images', 'grid'}
        or (layout == 'auto' and len(images) > 1 and spec.get('auto_multi', True))
    ):
        img_paths = _resolved_images(images)
        add_multi_image_slide(
            prs,
            label,
            img_paths,
            bullets=bullets_override,
            width_weights=spec.get('width_weights'),
            layout='grid' if layout == 'grid' else 'auto',
        )
        return

    slide = prs.slides.add_slide(_PRS_CTX['blank_layout'])
    add_header(slide, label)

    if images:
        layout_image_left(
            slide,
            images,
            others,
            bullets_override=bullets_override,
            text_heavy=text_heavy,
        )
    elif bullets_override is not None:
        layout_full_width(slide, [{'type': 'list', 'items': list(bullets_override)}])
    else:
        layout_full_width(slide, others)


# ------------------------------------------------------------------
# 5. main()
# ------------------------------------------------------------------
def main():
    prs = _load_template()
    _PRS_CTX['blank_layout'] = prs.slide_layouts[6]   # 空白
    title_layout              = prs.slide_layouts[0]   # 标题幻灯片
    chunks = load_chunks(fix_img_for_pages=False)      # PPTX 跟 REPORT.md 同根,不要改图路径

    # 5a. 标题页 —— 改这里
    title_slide = prs.slides.add_slide(title_layout)
    title_slide.shapes.title.text = '<项目名> × <子主题>'
    if len(title_slide.placeholders) > 1:
        sub = title_slide.placeholders[1]
        sub.text = '<作者> · <单位>\n<日期>'

    # 5b. content slides
    for spec in _iter_ppt_specs():
        slug = spec['slug']
        chunk_key = spec.get('chunk', slug)
        _PPT_CTX.clear()
        _PPT_CTX.update({'slug': slug, 'chunk': chunk_key})
        md_chunk = chunks.get(chunk_key)
        if md_chunk is None:
            _diag('missing_chunks', f'{slug}: chunk {chunk_key!r} not found')
            continue
        add_content_slide(prs, spec.get('label', slug), md_chunk, spec=spec)

    # 5c. 保存(被 PowerPoint 锁住时落到带时间戳的备用名)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(OUT))
        target = OUT
    except PermissionError:
        from datetime import datetime
        target = OUT.with_stem(OUT.stem + '_' + datetime.now().strftime('%H%M%S'))
        prs.save(str(target))
        print(f'WARN: {OUT} is locked (PowerPoint 还开着?), saved to alt path:')
    print(f'Wrote: {target}')
    print(f'  slides: {len(prs.slides)}')
    print(f'  size:   {target.stat().st_size:,} B')
    _print_diagnostics()


if __name__ == '__main__':
    main()
